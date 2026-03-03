import fitz
from pptx import Presentation


def extract_text_from_pdf(uploaded_file):
    text = ""
    pdf = fitz.open(stream=uploaded_file.read(), filetype="pdf")
    for page in pdf:
        text += page.get_text()
    return text


def extract_text_from_ppt(uploaded_file):
    text = ""
    prs = Presentation(uploaded_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text